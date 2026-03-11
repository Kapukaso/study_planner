"""Test document parsers."""
import os
import tempfile
from pathlib import Path

from src.parsers.pdf_parser import PDFParser
from src.parsers.docx_parser import DOCXParser
from src.parsers.ppt_parser import PPTParser
from src.parsers.image_parser import ImageParser
from PIL import Image, ImageDraw, ImageFont


def test_pdf_parser():
    """Test PDF parser with basic functionality."""
    parser = PDFParser()

    # Create a simple test PDF
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            pdf_path = tmp.name

        c = canvas.Canvas(pdf_path, pagesize=letter)
        c.setFont("Helvetica", 12)
        c.drawString(50, 750, "Test PDF content for parsing.")
        c.drawString(50, 720, "This is a second line of text.")
        c.save()

        chunks = parser.parse(pdf_path)
        print(f"PDF Parser: Found {len(chunks)} chunks")

        if chunks:
            print(f"First chunk text: {chunks[0].text[:100]}...")
            print(f"Metadata: {chunks[0].metadata}")

        # Test page count
        page_count = parser.get_page_count(pdf_path)
        print(f"PDF pages: {page_count}")

        # Cleanup
        os.unlink(pdf_path)
        return True

    except ImportError:
        print("PDF Parser: reportlab not available, skipping PDF creation test")
        return True
    except Exception as e:
        print(f"PDF Parser test failed: {e}")
        return False


def test_docx_parser():
    """Test DOCX parser with basic functionality."""
    parser = DOCXParser()

    try:
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            docx_path = tmp.name

        # Create a simple DOCX
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test paragraph with some content.')
        doc.add_paragraph('Another paragraph with different text.')

        # Add a simple table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = 'Header 1'
        table.cell(0, 1).text = 'Header 2'
        table.cell(1, 0).text = 'Data 1'
        table.cell(1, 1).text = 'Data 2'

        doc.save(docx_path)

        chunks = parser.parse(docx_path)
        print(f"DOCX Parser: Found {len(chunks)} chunks")

        if chunks:
            print(f"First chunk text: {chunks[0].text[:100]}...")
            print(f"Metadata: {chunks[0].metadata}")

        # Test heading extraction
        headings = parser.extract_headings(docx_path)
        print(f"Headings found: {headings}")

        # Cleanup
        os.unlink(docx_path)
        return True

    except ImportError:
        print("DOCX Parser: python-docx not available")
        return False
    except Exception as e:
        print(f"DOCX Parser test failed: {e}")
        return False


def test_ppt_parser():
    """Test PPT parser with basic functionality."""
    parser = PPTParser()

    try:
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            pptx_path = tmp.name

        # Create a simple PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Presentation"

        # Add content
        content = slide.placeholders[1]
        content.text = "This is slide content.\nSecond line of content."

        prs.save(pptx_path)

        chunks = parser.parse(pptx_path)
        print(f"PPT Parser: Found {len(chunks)} chunks")

        if chunks:
            print(f"First chunk text: {chunks[0].text[:100]}...")
            print(f"Metadata: {chunks[0].metadata}")

        # Test slide count
        slide_count = parser.get_slide_count(pptx_path)
        print(f"PPT slides: {slide_count}")

        # Cleanup
        os.unlink(pptx_path)
        return True

    except ImportError:
        print("PPT Parser: python-pptx not available")
        return False
    except Exception as e:
        print(f"PPT Parser test failed: {e}")
        return False


def test_image_parser():
    """Test Image parser with basic functionality."""
    parser = ImageParser()
    if not parser.tesseract_available:
        print("Image Parser: Tesseract not available, skipping test.")
        return True  # Skip test if OCR is not available

    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            image_path = tmp.name

        # Create a simple image with text
        img = Image.new('RGB', (400, 100), color='white')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 15)
        except IOError:
            font = ImageFont.load_default()
        
        test_text = "This is a test image for OCR."
        draw.text((10, 10), test_text, fill='black', font=font)
        img.save(image_path)

        # Parse the image
        chunks = parser.parse(image_path)
        print(f"Image Parser: Found {len(chunks)} chunks")

        assert len(chunks) > 0, "Parser should find at least one chunk."
        
        # The OCR might not be perfect, so we check for containment
        full_text = " ".join(c.text for c in chunks)
        print(f"Extracted text: {full_text}")
        # A simple check to see if most of the text was extracted
        assert "test image for OCR" in full_text, "Extracted text does not match expected content."

        if chunks:
            print(f"First chunk text: {chunks[0].text[:100]}...")
            print(f"Metadata: {chunks[0].metadata}")
            assert 'ocr_confidence' in chunks[0].metadata
            assert chunks[0].metadata['ocr_confidence'] > 0

        # Cleanup
        os.unlink(image_path)
        return True

    except Exception as e:
        print(f"Image Parser test failed: {e}")
        # Cleanup in case of failure
        if 'image_path' in locals() and os.path.exists(image_path):
            os.unlink(image_path)
        return False


if __name__ == "__main__":
    print("Testing Document Parsers")
    print("=" * 50)

    results = []

    print("\n1. Testing PDF Parser...")
    results.append(("PDF Parser", test_pdf_parser()))

    print("\n2. Testing DOCX Parser...")
    results.append(("DOCX Parser", test_docx_parser()))

    print("\n3. Testing PPT Parser...")
    results.append(("PPT Parser", test_ppt_parser()))

    print("\n4. Testing Image Parser...")
    results.append(("Image Parser", test_image_parser()))

    print("\n" + "=" * 50)
    print("Test Results:")
    for name, success in results:
        status = "PASS" if success else "FAIL"
        print(f"{name}: {status}")

    passed = sum(1 for _, success in results if success)
    print(f"\nPassed: {passed}/{len(results)} tests")