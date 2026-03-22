"""PPT/PPTX parser using python-pptx."""
from pptx import Presentation
from typing import List
import os

from src.parsers.base_parser import BaseParser, ParsedChunk
from src.parsers.exceptions import CorruptedFileError, ParserError


class PPTParser(BaseParser):
    """Parser for PPT/PPTX presentations."""
    
    def parse(self, file_path: str) -> List[ParsedChunk]:
        """
        Parse PPT/PPTX and extract text from slides.
        
        Args:
            file_path: Path to PPT/PPTX file
            
        Returns:
            List of ParsedChunk objects, one per slide
        """
        self.validate_file(file_path, ['.pptx', '.ppt'])
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                try:
                    # Extract text from all shapes in slide
                    slide_text_parts = []
                    shape_info = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                            shape_info.append({
                                'type': 'textbox',
                                'has_text': True
                            })
                        elif hasattr(shape, "image"):
                            # Note: images are present but we can't extract text from them
                            shape_info.append({
                                'type': 'image',
                                'has_text': False
                            })
                        else:
                            shape_info.append({
                                'type': 'other',
                                'has_text': False
                            })
                    
                    # Combine slide text
                    if slide_text_parts:
                        text = '\n\n'.join(slide_text_parts)
                        cleaned_text = self.clean_text(text)
                        
                        if cleaned_text:
                            chunk = ParsedChunk(
                                text=cleaned_text,
                                page_number=slide_num,  # Treat slide as "page"
                                chunk_index=slide_num - 1,
                                metadata={
                                    'slide_number': slide_num,
                                    'total_shapes': len(slide.shapes),
                                    'text_shapes': len([s for s in shape_info if s['has_text']]),
                                    'image_shapes': len([s for s in shape_info if s['type'] == 'image']),
                                    'parser': 'python-pptx',
                                    'slide_title': slide.shapes.title.text if slide.shapes.title else None
                                }
                            )
                            chunks.append(chunk)
                except Exception as slide_err:
                    print(f"Warning: Failed to parse slide {slide_num} of {file_path}: {str(slide_err)}")
                    continue
        
        except Exception as e:
            if "Package not found" in str(e) or "not a PowerPoint file" in str(e):
                raise CorruptedFileError(f"PPT file is corrupted or invalid: {file_path}")
            if isinstance(e, ParserError):
                raise e
            raise ParserError(f"Error parsing PPT: {str(e)}")
        
        return chunks
    
    def get_slide_count(self, file_path: str) -> int:
        """Get total number of slides."""
        try:
            prs = Presentation(file_path)
            return len(prs.slides)
        except Exception:
            return 0
